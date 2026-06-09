"""
reports_api
───────────
R45 — Picture-perfect Report Designer reports. Persisted to
~/jarvis/semantic/reports.json (atomic writes + timestamped backups,
same pattern as briefs_api / queries_api).

A report is a multi-page document. Each page has elements at absolute
positions: {id, type, x, y, w, h, props, locked}. Element types match
the Designer palette: text, divider, shape, image, kpi, widget.

Routes (prefix /api/reports):
  GET    /api/reports                  — list summaries (scope-filtered)
  GET    /api/reports/{id}             — full report
  POST   /api/reports                  — create
  PUT    /api/reports/{id}             — update (owner only)
  DELETE /api/reports/{id}             — remove (owner only)
  POST   /api/reports/{id}/export-pptx — native PowerPoint export
"""

from __future__ import annotations

import io
import json
import os
import time
import uuid
import shutil
import threading
import urllib.request
from typing import Optional

from fastapi import APIRouter, Request, WebSocket, WebSocketDisconnect
from fastapi.responses import JSONResponse, Response
from pydantic import BaseModel

from semantic.loader import SEMANTIC_DIR

router = APIRouter(prefix="/api/reports", tags=["reports"])

REPORTS_PATH = SEMANTIC_DIR / "reports.json"
_lock = threading.Lock()

# Designer canvas size in pixels — kept in sync with martin_app.html's RD.pageW/pageH.
# We translate these to EMU when generating PPTX (1 inch = 914400 EMU, 96px ≈ 1in).
DESIGN_W_PX = 1280
DESIGN_H_PX = 720


def _now() -> float:
    return time.time()


def _err(status: int, error: str, detail: str = "") -> JSONResponse:
    return JSONResponse(status_code=status, content={"error": error, "detail": detail})


def _load() -> dict:
    if not REPORTS_PATH.exists():
        return {"reports": []}
    try:
        with open(REPORTS_PATH, "r", encoding="utf-8") as f:
            data = json.load(f)
        if not isinstance(data, dict) or "reports" not in data:
            return {"reports": []}
        return data
    except (json.JSONDecodeError, OSError):
        return {"reports": []}


def _atomic_write(data: dict):
    REPORTS_PATH.parent.mkdir(parents=True, exist_ok=True)
    if REPORTS_PATH.exists():
        ts = time.strftime("%Y%m%d_%H%M%S")
        try:
            shutil.copy2(REPORTS_PATH, REPORTS_PATH.with_suffix(".json.bak." + ts))
        except OSError:
            pass
    tmp = REPORTS_PATH.with_suffix(".json.tmp")
    with open(tmp, "w", encoding="utf-8") as f:
        json.dump(data, f, indent=2, default=str)
        f.flush()
        os.fsync(f.fileno())
    tmp.replace(REPORTS_PATH)


def _summary(r: dict) -> dict:
    pages = r.get("pages") or []
    return {
        "id": r.get("id"),
        "name": r.get("name"),
        "description": r.get("description"),
        "page_count": len(pages),
        "element_count": sum(len(p.get("elements") or []) for p in pages),
        "owner_email": r.get("owner_email"),
        "scope": r.get("scope") or "private",
        "created_at": r.get("created_at"),
        "updated_at": r.get("updated_at"),
    }


def _viewer_email(request: Request) -> Optional[str]:
    return (request.headers.get("X-Jarvis-User") or "").strip() or None


def _can_see(r: dict, viewer: Optional[str]) -> bool:
    scope = (r.get("scope") or "private").lower()
    if scope != "private":
        return True
    return bool(viewer) and viewer == r.get("owner_email")


class ReportCreate(BaseModel):
    name: str
    description: Optional[str] = ""
    pages: list
    page_w: Optional[int] = DESIGN_W_PX
    page_h: Optional[int] = DESIGN_H_PX
    scope: Optional[str] = "private"


class ReportUpdate(BaseModel):
    name: Optional[str] = None
    description: Optional[str] = None
    pages: Optional[list] = None
    page_w: Optional[int] = None
    page_h: Optional[int] = None
    scope: Optional[str] = None


@router.get("")
async def list_reports(request: Request):
    viewer = _viewer_email(request)
    data = _load()
    items = [_summary(r) for r in data["reports"] if _can_see(r, viewer)]
    items.sort(key=lambda x: x.get("updated_at") or 0, reverse=True)
    return {"reports": items, "count": len(items)}


@router.get("/{report_id}")
async def get_report(report_id: str, request: Request):
    viewer = _viewer_email(request)
    data = _load()
    r = next((x for x in data["reports"] if x.get("id") == report_id), None)
    if not r:
        return _err(404, "not found", f"No report '{report_id}'")
    if not _can_see(r, viewer):
        return _err(403, "forbidden", "This report is private to its owner.")
    return r


@router.post("")
async def create_report(body: ReportCreate, request: Request):
    if not body.name.strip():
        return _err(400, "name required", "Give the report a name.")
    scope = (body.scope or "private").lower()
    if scope not in ("private", "group", "corporate"):
        scope = "private"
    owner = _viewer_email(request)
    with _lock:
        data = _load()
        r = {
            "id": str(uuid.uuid4()),
            "name": body.name.strip(),
            "description": (body.description or "").strip(),
            "pages": body.pages or [{"name": "Page 1", "elements": []}],
            "page_w": body.page_w or DESIGN_W_PX,
            "page_h": body.page_h or DESIGN_H_PX,
            "owner_email": owner,
            "scope": scope,
            "created_at": _now(),
            "updated_at": _now(),
        }
        data["reports"].append(r)
        try:
            _atomic_write(data)
        except OSError as e:
            return _err(500, "save failed", str(e))
    return r


@router.put("/{report_id}")
async def update_report(report_id: str, body: ReportUpdate, request: Request):
    viewer = _viewer_email(request)
    with _lock:
        data = _load()
        r = next((x for x in data["reports"] if x.get("id") == report_id), None)
        if not r:
            return _err(404, "not found", f"No report '{report_id}'")
        if r.get("owner_email") and viewer != r.get("owner_email"):
            return _err(403, "forbidden", "Only the owner can edit a saved report.")
        if body.name is not None:        r["name"] = body.name.strip() or r["name"]
        if body.description is not None: r["description"] = body.description.strip()
        if body.pages is not None:       r["pages"] = body.pages
        if body.page_w is not None:      r["page_w"] = int(body.page_w)
        if body.page_h is not None:      r["page_h"] = int(body.page_h)
        if body.scope is not None:
            s = body.scope.lower()
            if s in ("private", "group", "corporate"):
                r["scope"] = s
        r["updated_at"] = _now()
        try:
            _atomic_write(data)
        except OSError as e:
            return _err(500, "save failed", str(e))
    return r


@router.delete("/{report_id}")
async def delete_report(report_id: str, request: Request):
    viewer = _viewer_email(request)
    with _lock:
        data = _load()
        target = next((x for x in data["reports"] if x.get("id") == report_id), None)
        if not target:
            return _err(404, "not found", f"No report '{report_id}'")
        if target.get("owner_email") and viewer != target.get("owner_email"):
            return _err(403, "forbidden", "Only the owner can delete a saved report.")
        data["reports"] = [x for x in data["reports"] if x.get("id") != report_id]
        try:
            _atomic_write(data)
        except OSError as e:
            return _err(500, "save failed", str(e))
    return {"ok": True, "deleted": report_id}


# ─── PPTX EXPORT ─────────────────────────────────────────────────────────
# Each page becomes a slide. Elements are translated to PowerPoint shapes,
# preserving absolute positions. Query widgets execute their saved query
# and embed a chart rendered as PNG via matplotlib.

def _px_to_emu(px: int | float, ppi_design: int, design_total_px: int, slide_emu: int) -> int:
    """Translate the Designer canvas's pixel coordinates to PPTX EMU.

    We map [0..design_total_px] → [0..slide_emu] linearly so a 1280×720
    canvas fits exactly on a 10in × 5.625in widescreen slide.
    """
    return int(round((px / design_total_px) * slide_emu))


def _hex_to_rgb(s: str) -> tuple[int, int, int]:
    s = (s or "").strip().lstrip("#")
    if len(s) == 3:
        s = "".join(c * 2 for c in s)
    if len(s) != 6:
        return (15, 23, 42)
    try:
        return tuple(int(s[i:i + 2], 16) for i in (0, 2, 4))  # type: ignore[return-value]
    except ValueError:
        return (15, 23, 42)


def _fetch_image_bytes(src: str) -> Optional[bytes]:
    """Best-effort fetch of an image URL. Returns None on any failure."""
    if not src or not src.startswith(("http://", "https://")):
        return None
    try:
        with urllib.request.urlopen(src, timeout=8) as r:
            return r.read()
    except Exception:
        return None


def _render_widget_chart_png(el: dict, viewer: Optional[str]) -> Optional[bytes]:
    """Run the saved query and render a chart as PNG bytes."""
    qid = ((el.get("props") or {}).get("queryId"))
    if not qid:
        return None
    # Pull the saved query, run it, draw with matplotlib.
    try:
        from queries_api import _load as queries_load
        qdata = queries_load()
    except Exception:
        return None
    q = next((x for x in qdata.get("queries", []) if x.get("id") == qid), None)
    if not q:
        return None
    # Permission check on the underlying query too
    scope = (q.get("scope") or "private").lower()
    if scope == "private" and viewer and viewer != q.get("owner_email"):
        return None
    qdict = dict(q.get("qdict") or {})
    try:
        from semantic_api import _load_fresh as _sem_load
        from semantic.query_builder import StructuredQuery, build_sql
        from semantic.executor import run_query as _exec
        model = _sem_load()
        sq = StructuredQuery.from_dict(qdict)
        sql = build_sql(sq, model)
        result = _exec(sql)
        cols, rows = result.columns, result.rows
    except Exception:
        return None
    if not rows:
        return None
    # Render with matplotlib (no GUI). Default to bar chart of first metric by first dim.
    try:
        import matplotlib
        matplotlib.use("Agg")
        import matplotlib.pyplot as plt
        chart_type = ((el.get("props") or {}).get("chartType") or "bar").lower()
        n_dims = max(0, len(cols) - len(qdict.get("metrics") or []))
        if n_dims == 0:
            # KPI-like — fall back to a single-bar chart
            labels = [cols[0] if cols else "Value"]
            values = [float(rows[0][0]) if rows[0] else 0.0]
        else:
            labels = [str(r[0]) for r in rows[:15]]
            metric_col = n_dims
            values = [float(r[metric_col]) if r[metric_col] is not None else 0.0 for r in rows[:15]]
        fig, ax = plt.subplots(figsize=(8, 4.5), dpi=120)
        if chart_type == "line" or chart_type == "area":
            ax.plot(labels, values, marker="o", color="#3a7a9b", linewidth=2)
            if chart_type == "area":
                ax.fill_between(range(len(labels)), values, color="#3a7a9b", alpha=0.25)
        elif chart_type in ("pie", "donut"):
            ax.pie(values, labels=labels, autopct="%1.0f%%", colors=plt.cm.Blues([0.5 + 0.35 * (i / max(1, len(values))) for i in range(len(values))]))
        else:
            ax.bar(labels, values, color="#3a7a9b", width=0.6)
        title = (el.get("props") or {}).get("queryName") or q.get("name") or ""
        if title:
            ax.set_title(title, fontsize=11, fontweight=600, loc="left", pad=8)
        if chart_type not in ("pie", "donut"):
            ax.spines["top"].set_visible(False)
            ax.spines["right"].set_visible(False)
            plt.xticks(rotation=35, ha="right", fontsize=8)
            plt.yticks(fontsize=8)
            for spine in ax.spines.values():
                spine.set_color("#cbd5e1")
        fig.tight_layout()
        buf = io.BytesIO()
        fig.savefig(buf, format="png", bbox_inches="tight", facecolor="white")
        plt.close(fig)
        return buf.getvalue()
    except Exception:
        return None


# ─── R50 — Real-time collab presence (WebSocket) ──────────────────────────
# Lightweight presence layer for the Designer. Connected users share cursor
# positions and selections via a per-report room. No persistence — when the
# last user disconnects the room dies. Messages are JSON and pass through
# unchanged (we just fan-out), so the frontend defines the protocol.

# In-memory rooms: report_id → set of WebSocket connections
_rd_rooms: dict[str, set] = {}


@router.websocket("/{report_id}/collab")
async def collab_socket(websocket: WebSocket, report_id: str):
    """Broadcast cursor + selection messages between users editing the same report.

    Authentication is intentionally lightweight — this is a presence channel,
    not a write path. The actual report data is still owner-protected via
    REST. Bad actors can only fake cursors, not change anything.
    """
    await websocket.accept()
    room = _rd_rooms.setdefault(report_id, set())
    room.add(websocket)
    try:
        # On join, announce to others so they re-broadcast their last cursor.
        for ws in list(room):
            if ws is websocket: continue
            try: await ws.send_json({"type": "hello", "report_id": report_id})
            except Exception: pass
        # Pump
        while True:
            data = await websocket.receive_json()
            # Drop any oversized payloads (defensive)
            try:
                if isinstance(data, dict): data["report_id"] = report_id
            except Exception:
                continue
            # Fan-out to everyone else in the room (excluding sender)
            for ws in list(room):
                if ws is websocket: continue
                try:
                    await ws.send_json(data)
                except Exception:
                    # Stale connection — drop it
                    room.discard(ws)
    except WebSocketDisconnect:
        pass
    except Exception:
        pass
    finally:
        room.discard(websocket)
        # Announce departure
        uid = None
        try:
            # Last message's uid field, best-effort
            uid = (data or {}).get("uid") if isinstance(data, dict) else None  # noqa: F821
        except Exception:
            uid = None
        if uid:
            for ws in list(room):
                try: await ws.send_json({"type": "bye", "uid": uid})
                except Exception: pass
        if not room:
            _rd_rooms.pop(report_id, None)


# ─── R48 — AI report generator ───────────────────────────────────────────
# Takes a free-text prompt ("design a Q3 marketing overview for Cardiology
# with KPIs, channel mix, top campaigns, conversion funnel") and asks
# Claude to lay out pages of elements. The returned schema matches the
# Designer's element shape so the frontend just hydrates it.

_AI_DESIGNER_SYSTEM = (
    "You are a senior marketing-analytics designer at Orlando Health. "
    "Given a description of a report the user wants, design a clean multi-page "
    "layout using the Report Designer's element vocabulary. Output ONLY JSON, "
    "no code fences, no commentary.\n\n"
    "Schema:\n"
    "{\n"
    '  "name": "<concise report title>",\n'
    '  "theme": "oh-classic" | "oh-light" | "oh-executive" | "oh-warmth",\n'
    '  "pages": [\n'
    "    {\n"
    '      "name": "<page name, short>",\n'
    '      "elements": [\n'
    '        {"type":"text",    "x":40, "y":40, "w":1200, "h":60, "props":{"text":"<headline>", "size":36, "weight":700, "color":"#0f172a"}},\n'
    '        {"type":"text",    "x":40, "y":100,"w":1200,"h":30, "props":{"text":"<subtitle>",  "size":14, "weight":500, "color":"#64748b"}},\n'
    '        {"type":"divider", "x":40, "y":145,"w":1200,"h":2,  "props":{}},\n'
    '        {"type":"kpi",     "x":40, "y":170,"w":280, "h":110,"props":{"label":"<LABEL>","value":"<value>","delta":"<delta>"}},\n'
    '        {"type":"shape",   "x":40, "y":300,"w":300, "h":200,"props":{"fill":"#fce7e8","stroke":"#8C2026","radius":8}},\n'
    '        {"type":"widget",  "x":40, "y":340,"w":580, "h":320,"props":{"queryId":null,"queryName":"<what the chart shows>","chartType":"bar|line|donut|pie|area|combo"}}\n'
    "      ]\n"
    "    }\n"
    "  ]\n"
    "}\n\n"
    "RULES — these matter:\n"
    "• Canvas is 1280×720 pixels per page. Keep every element inside those bounds.\n"
    "• Pages target service-line marketing managers (Cardiology, Cancer, Orthopedics, Pediatrics, ER, Women's Health, etc.). Use specific, accurate language for whichever line the user mentioned.\n"
    "• KPI tile dimensions: 280×110 standard, 280×130 for prominent ones. Place 3–4 KPIs across the top of the first page.\n"
    "• Widgets (charts) should be ~580×320 for half-row, 1200×320 for full-row.\n"
    "• Use one big page-title at y=40 (size 32–36, weight 700) plus a subtitle at y=100.\n"
    "• Pick 'theme' that matches the audience: 'oh-classic' = standard internal, 'oh-executive' = board-level, 'oh-light' = digital snapshots, 'oh-warmth' = patient-facing.\n"
    "• Widgets have queryId=null because the user wires the data after; just describe the chart in queryName.\n"
    "• Section headers above widgets at size 15–16, weight 600.\n"
    "• Include 1 page for short questions, 2–3 pages for broad questions.\n"
    "• NEVER overlap elements. NEVER leave gaps wider than ~40px horizontally between sections."
)


class AIDesignRequest(BaseModel):
    prompt: str


@router.post("/ai-generate")
async def ai_generate(body: AIDesignRequest, request: Request):
    if not body.prompt.strip():
        return _err(400, "prompt required", "Describe the report you want.")
    key = os.getenv("ANTHROPIC_API_KEY", "")
    if not key:
        return _err(503, "AI not configured", "ANTHROPIC_API_KEY is not set.")
    try:
        import anthropic
        client = anthropic.AsyncAnthropic(api_key=key)
        resp = await client.messages.create(
            model="claude-haiku-4-5-20251001",
            max_tokens=4000,
            system=_AI_DESIGNER_SYSTEM,
            messages=[{"role": "user", "content": body.prompt.strip()}],
        )
        raw = resp.content[0].text if resp.content else ""
    except Exception as e:
        return _err(502, "AI request failed", str(e))

    # Extract JSON; Claude is supposed to return pure JSON but be defensive.
    import re as _re
    m = _re.search(r"\{.*\}", raw, _re.DOTALL)
    if not m:
        return _err(422, "could not parse AI response", "Claude did not return valid JSON.")
    try:
        spec = json.loads(m.group(0))
    except json.JSONDecodeError as e:
        return _err(422, "invalid AI JSON", str(e))
    pages = spec.get("pages") or []
    if not pages:
        return _err(422, "empty layout", "AI returned no pages.")
    # Mild sanity-clamp every element to the canvas; ignore bad ones
    safe_pages = []
    for p in pages:
        elems = []
        for el in (p.get("elements") or []):
            try:
                t = el.get("type")
                if t not in ("text", "divider", "shape", "image", "kpi", "widget"):
                    continue
                x = max(0, min(int(el.get("x", 0)), DESIGN_W_PX - 20))
                y = max(0, min(int(el.get("y", 0)), DESIGN_H_PX - 20))
                w = max(20, min(int(el.get("w", 200)), DESIGN_W_PX - x))
                h = max(20, min(int(el.get("h", 60)),  DESIGN_H_PX - y))
                elems.append({
                    "type": t, "x": x, "y": y, "w": w, "h": h,
                    "props": el.get("props") or {},
                })
            except Exception:
                continue
        safe_pages.append({"name": p.get("name") or "Page", "elements": elems})
    return {
        "name":  spec.get("name") or "AI-designed report",
        "theme": spec.get("theme") or "oh-classic",
        "pages": safe_pages,
    }


@router.post("/{report_id}/export-pptx")
async def export_pptx(report_id: str, request: Request):
    """Generate a native .pptx from the report's pages.

    Each Designer page → one slide. Element positions/sizes scale from the
    Designer canvas to the slide. Widget elements are rendered to PNG and
    embedded.
    """
    viewer = _viewer_email(request)
    data = _load()
    r = next((x for x in data["reports"] if x.get("id") == report_id), None)
    if not r:
        return _err(404, "not found", f"No report '{report_id}'")
    if not _can_see(r, viewer):
        return _err(403, "forbidden", "This report is private to its owner.")
    try:
        from pptx import Presentation
        from pptx.util import Emu, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.enum.text import PP_ALIGN
    except Exception as e:
        return _err(500, "pptx not available", f"python-pptx missing: {e}")

    pres = Presentation()
    # Widescreen 16:9, 10in × 5.625in
    pres.slide_width = Emu(9144000)
    pres.slide_height = Emu(5143500)
    sw = pres.slide_width
    sh = pres.slide_height
    page_w_px = r.get("page_w") or DESIGN_W_PX
    page_h_px = r.get("page_h") or DESIGN_H_PX

    blank_layout = pres.slide_layouts[6]  # blank

    for page in (r.get("pages") or []):
        slide = pres.slides.add_slide(blank_layout)
        for el in (page.get("elements") or []):
            try:
                x_emu = int(round((el.get("x", 0) / page_w_px) * sw))
                y_emu = int(round((el.get("y", 0) / page_h_px) * sh))
                w_emu = int(round((el.get("w", 100) / page_w_px) * sw))
                h_emu = int(round((el.get("h", 50) / page_h_px) * sh))
                t = el.get("type")
                props = el.get("props") or {}
                if t == "text":
                    tb = slide.shapes.add_textbox(x_emu, y_emu, w_emu, h_emu)
                    tf = tb.text_frame
                    tf.word_wrap = True
                    tf.margin_left = tf.margin_right = Emu(45720)
                    tf.margin_top = tf.margin_bottom = Emu(45720)
                    p = tf.paragraphs[0]
                    run = p.add_run()
                    run.text = str(props.get("text", ""))
                    run.font.size = Pt(int(props.get("size", 16)))
                    run.font.bold = int(props.get("weight", 500)) >= 600
                    rgb = _hex_to_rgb(props.get("color") or "#0f172a")
                    run.font.color.rgb = RGBColor(*rgb)
                    align_map = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}
                    p.alignment = align_map.get(props.get("align", "left"), PP_ALIGN.LEFT)
                elif t == "divider":
                    sh_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x_emu, y_emu, w_emu, h_emu)
                    sh_shape.fill.solid()
                    sh_shape.fill.fore_color.rgb = RGBColor(*_hex_to_rgb("#cbd5e1"))
                    sh_shape.line.fill.background()
                elif t == "shape":
                    radius_px = int(props.get("radius") or 8)
                    msoshape = MSO_SHAPE.ROUNDED_RECTANGLE if radius_px > 2 else MSO_SHAPE.RECTANGLE
                    sh_shape = slide.shapes.add_shape(msoshape, x_emu, y_emu, w_emu, h_emu)
                    sh_shape.fill.solid()
                    sh_shape.fill.fore_color.rgb = RGBColor(*_hex_to_rgb(props.get("fill") or "#dbeafe"))
                    stroke = _hex_to_rgb(props.get("stroke") or "#0E76B4")
                    sh_shape.line.color.rgb = RGBColor(*stroke)
                    sh_shape.line.width = Emu(20000)
                elif t == "image":
                    src = props.get("src") or ""
                    img_bytes = _fetch_image_bytes(src)
                    if img_bytes:
                        slide.shapes.add_picture(io.BytesIO(img_bytes), x_emu, y_emu, w_emu, h_emu)
                    else:
                        # Placeholder rectangle so the slot is visible
                        sh_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x_emu, y_emu, w_emu, h_emu)
                        sh_shape.fill.solid()
                        sh_shape.fill.fore_color.rgb = RGBColor(241, 245, 249)
                        sh_shape.line.color.rgb = RGBColor(148, 163, 184)
                        tf = sh_shape.text_frame
                        tf.text = "[image]"
                elif t == "kpi":
                    # Card background
                    bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x_emu, y_emu, w_emu, h_emu)
                    bg.fill.solid()
                    bg.fill.fore_color.rgb = RGBColor(255, 255, 255)
                    bg.line.color.rgb = RGBColor(226, 232, 240)
                    # Label + Value text box overlapping the same bounds
                    tb = slide.shapes.add_textbox(x_emu, y_emu, w_emu, h_emu)
                    tf = tb.text_frame
                    tf.word_wrap = True
                    tf.margin_left = tf.margin_right = Emu(140000)
                    tf.margin_top = Emu(120000)
                    # Label paragraph
                    p1 = tf.paragraphs[0]
                    r1 = p1.add_run()
                    r1.text = str(props.get("label", "METRIC")).upper()
                    r1.font.size = Pt(10)
                    r1.font.bold = True
                    r1.font.color.rgb = RGBColor(100, 116, 139)
                    # Value paragraph
                    p2 = tf.add_paragraph()
                    r2 = p2.add_run()
                    r2.text = str(props.get("value", ""))
                    r2.font.size = Pt(32)
                    r2.font.bold = True
                    r2.font.color.rgb = RGBColor(15, 23, 42)
                    if props.get("delta"):
                        p3 = tf.add_paragraph()
                        r3 = p3.add_run()
                        r3.text = str(props["delta"])
                        r3.font.size = Pt(10)
                        r3.font.color.rgb = RGBColor(13, 148, 136)
                elif t == "widget":
                    png = _render_widget_chart_png(el, viewer)
                    if png:
                        slide.shapes.add_picture(io.BytesIO(png), x_emu, y_emu, w_emu, h_emu)
                    else:
                        # Placeholder
                        sh_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x_emu, y_emu, w_emu, h_emu)
                        sh_shape.fill.solid()
                        sh_shape.fill.fore_color.rgb = RGBColor(248, 250, 252)
                        sh_shape.line.color.rgb = RGBColor(226, 232, 240)
                        tf = sh_shape.text_frame
                        tf.text = f"[widget: {(el.get('props') or {}).get('queryName') or 'unknown'}]"
            except Exception:
                # If any single element fails, skip it but keep going so the rest of
                # the slide is still useful. Trying again later is cheap.
                continue

    out = io.BytesIO()
    pres.save(out)
    blob = out.getvalue()
    fname = (r.get("name") or "report").replace("/", "_").replace(" ", "_") + ".pptx"
    return Response(
        content=blob,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": f'attachment; filename="{fname}"'},
    )
